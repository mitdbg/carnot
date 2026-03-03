"""File content parsing utilities.

Extracts the file-reading logic previously embedded in ``DataItem._read_file_contents()``
into a standalone function that works with raw bytes, so the storage layer can
read the bytes and this module can parse them.
"""

from __future__ import annotations

import io


def parse_file_contents(uri: str, raw: bytes) -> str:
    """Parse raw bytes from a file into a text string based on file extension.

    Parameters
    ----------
    uri:
        The file path or URI.  Only the extension (suffix) is used to
        determine the parser.
    raw:
        The raw bytes of the file.

    Returns
    -------
    str
        The parsed textual content.
    """
    suffix = uri.rsplit(".", 1)[-1].lower() if "." in uri else ""

    try:
        if suffix == "txt":
            return raw.decode("utf-8")

        elif suffix == "csv":
            import pandas as pd
            df = pd.read_csv(io.BytesIO(raw))
            return df.to_string()

        elif suffix == "parquet":
            import pandas as pd
            df = pd.read_parquet(io.BytesIO(raw))
            return df.to_string()

        elif suffix in ("xlsx", "xls"):
            import pandas as pd
            df = pd.read_excel(io.BytesIO(raw))
            return df.to_string()

        elif suffix == "pdf":
            import fitz
            doc = fitz.open(stream=raw, filetype="pdf")
            return "\n".join(page.get_text() for page in doc)

        elif suffix == "docx":
            from docx import Document
            doc = Document(io.BytesIO(raw))
            return "\n".join(para.text for para in doc.paragraphs)

        elif suffix == "pptx":
            from pptx import Presentation
            pres = Presentation(io.BytesIO(raw))
            text_runs: list[str] = []
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)

        elif suffix == "json":
            return raw.decode("utf-8")

        else:
            # Fallback: decode as UTF-8 with lossy handling
            return raw.decode("utf-8", errors="ignore")

    except Exception as e:
        return f"Error reading file {uri}: {e!s}"
